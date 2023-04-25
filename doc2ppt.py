""".doxc document to .ppt"""
import docx
from pptx import Presentation
from pptx.util import Inches

# Read Word document
def read_word_doc(filename):
    doc = docx.Document(filename)
    data = []

    for para in doc.paragraphs:
        data.append(para.text)
    
    return data

# create a slide with title and content
def create_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.text = content

# the main function to run the program
def main():
    word_file = 'doc2ppt_file.docx'
    ppt_file = 'doc2_ppt_slides.pptx'

    data = read_word_doc(word_file)
    prs = Presentation()

    for index, item in enumerate(data):
        if index % 2 == 0:
            title = item
        else:
            content = item
            create_slide(prs, title, content)
    
    prs.save(ppt_file)

if __name__ == "__main__":
    main()
